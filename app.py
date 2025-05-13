import openai
import stripe
import os
import io
import tempfile
import base64
import pytesseract
import toml
import pandas as pd
import json
from pptx import Presentation
from PIL import Image
import fitz  # PyMuPDF
import docx
from werkzeug.utils import secure_filename
import uuid
from flask import Flask, render_template, request, jsonify, send_file
from io import BytesIO
from langdetect import detect, DetectorFactory
from PIL import Image
import pytesseract
import requests



DetectorFactory.seed = 0  # agar hasil konsisten

def detect_language(text):
    try:
        lang = detect(text)
        return lang
    except:
        return "unknown"


def extract_images_from_pdf(stream):
    stream.seek(0)
    images = []
    with fitz.open(stream=stream.read(), filetype="pdf") as doc:
        for page in doc:
            for img in page.get_images(full=True):
                xref = img[0]
                base_image = doc.extract_image(xref)
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                image_bytes = base_image["image"]
                image = Image.open(BytesIO(image_bytes)).convert("RGB")
                images.append(image)
    return images

def load_search_history():
    if os.path.exists(history_file):
        try:
            with open(history_file, "r", encoding="utf-8") as f:
                history = json.load(f)
                return [entry for entry in history if "type" in entry and "input" in entry and "response" in entry]
        except json.JSONDecodeError:
            st.warning("⚠️ File riwayat pencarian rusak. Menghapus riwayat lama.")
            return []
    return []

HISTORY_FILE = "search_history.json"

def load_chat_history():
    if os.path.exists(HISTORY_FILE):
        with open(HISTORY_FILE, "r") as f:
            return json.load(f)
    return []

def save_chat_history(history):
    with open(HISTORY_FILE, "w") as f:
        json.dump(history, f, indent=2)

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_images_from_pptx(file):
    file.seek(0)
    prs = Presentation(file)
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image") and shape.image:
                image_stream = shape.image.blob
                image = Image.open(BytesIO(image_stream)).convert("RGB")
                images.append(image)
    return images

def extract_text_from_xlsx(file):
    file.seek(0)
    df = pd.read_excel(file, sheet_name=None)
    output = ""
    for sheet_name, sheet in df.items():
        output += f"\n--- Sheet: {sheet_name} ---\n"
        output += sheet.to_string(index=False)
    return output

def extract_text_from_image(file):
    image = Image.open(file).convert("RGB")
    return pytesseract.image_to_string(image)

def extract_text_from_images(images):
    return "\n\n".join([pytesseract.image_to_string(img) for img in images])

app = Flask(__name__)
app.secret_key = os.urandom(24)

# Load secrets
# Load secrets
secrets = toml.load("secrets.toml")
openai_api_key = secrets.get("OPENAI_API_KEY")
stripe.api_key = secrets.get("STRIPE_SECRET_KEY")
# Set API key untuk OpenAI
openai.api_key = openai_api_key  # <<< tambahkan baris ini
PREMIUM_FILE = 'premium_emails.json'

def load_premium_emails():
    if os.path.exists(PREMIUM_FILE):
        with open(PREMIUM_FILE, 'r') as f:
            return set(json.load(f))
    return set()

def save_premium_email(email):
    emails = load_premium_emails()
    emails.add(email)
    with open(PREMIUM_FILE, 'w') as f:
        json.dump(list(emails), f)

# Function to handle file extraction
def extract_text_from_images(images):
    text = ""
    for img in images:
        text += pytesseract.image_to_string(img) + "\n"
    return text


def extract_images_from_docx(file):
    images = []
    doc = docx.Document(file)
    for rel in doc.part._rels.values():
        if "image" in rel.target_ref:
            data = rel.target_part.blob
            images.append(Image.open(io.BytesIO(data)))
    return images

def extract_images_from_pptx(file):
    images = []
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                image_stream = shape.image.blob
                images.append(Image.open(io.BytesIO(image_stream)))
    return images

# Function to create Stripe checkout session
@app.route('/create-checkout-session', methods=['POST'])
def create_checkout_session():
    data = request.get_json()
    email = data.get("email")

    if not email:
        return jsonify({"error": "Email tidak valid."}), 400

    try:
        checkout_session = stripe.checkout.Session.create(
            payment_method_types=['card'],
            line_items=[{
                'price_data': {
                    'currency': 'usd',
                    'product_data': {'name': 'Akses Premium Doragonnoid Chatbot'},
                    'unit_amount': 1000,  # harga dalam cents (misalnya 10.00 USD)
                },
                'quantity': 1,
            }],
            mode='payment',
            success_url='http://localhost:5000/success?session_id={CHECKOUT_SESSION_ID}',
            cancel_url='http://localhost:5000/cancel',
        )

        return jsonify({'url': checkout_session.url})

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Route for payment success
@app.route('/success')
def success():
    email = request.args.get('email')
    if email:
        save_premium_email(email)
    return "Pembayaran sukses. Akses premium diberikan."

# File upload and text extraction route


# Route to check if email is premium
@app.route('/verify-email', methods=['POST'])
def verify_email():
    data = request.get_json()
    email = data.get('email')
    emails = load_premium_emails()
    return jsonify({"premium": email in emails})
def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_pdf(file_stream):
    text = ""
    with fitz.open(stream=file_stream.read(), filetype="pdf") as doc:
        for page in doc:
            text += page.get_text() + "\n"
    return text


@app.route('/upload', methods=['POST'])
def upload():
    file = request.files.get('file')
    if not file:
        return jsonify({'error': 'Tidak ada file diunggah'}), 400

    filename = file.filename.lower()
    try:
        text, ocr_text = "", ""

        if filename.endswith(".pdf"):
            file.stream.seek(0)
            text = extract_text_from_pdf(file.stream)
            file.stream.seek(0)
            ocr_text = extract_text_from_images(extract_images_from_pdf(file.stream))

        elif filename.endswith(".docx"):
            text = extract_text_from_docx(file)
            ocr_text = extract_text_from_images(extract_images_from_docx(file))

        elif filename.endswith(".pptx"):
            text = extract_text_from_pptx(file)
            ocr_text = extract_text_from_images(extract_images_from_pptx(file))

        elif filename.endswith(".xlsx"):
            text = extract_text_from_xlsx(file)

        elif filename.endswith((".png", ".jpg", ".jpeg")):
            text = extract_text_from_image(file)

        else:
            return jsonify({'error': 'Format tidak didukung'}), 400

        combined = (text + "\n\n[OCR Gambar]\n" + ocr_text).strip()
        detected_lang = detect_language(combined[:500])  # hanya pakai 500 karakter pertama untuk efisiensi

        return jsonify({
            'text': combined or "Tidak ditemukan teks.",
            'language': detected_lang
        })


    except Exception as e:
        return jsonify({'error': str(e)}), 500
    
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)  # Ini akan membuat folder jika belum ada
@app.route("/analyze-image", methods=["POST"])
def analyze_image():
    try:
        file = request.files["image"]
        instruction = request.form.get("instruction", "")
        filename = secure_filename(file.filename)
        unique_filename = f"{uuid.uuid4().hex}_{filename}"
        filepath = os.path.join("uploads", unique_filename)
        file.save(filepath)

        # Lakukan analisis dengan OpenAI Vision
        with open(filepath, "rb") as img_file:
            response = openai.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "user", "content": [
                        {"type": "text", "text": instruction or "Analisis gambar ini"},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64.b64encode(img_file.read()).decode()}"}}
                    ]}
                ],
                max_tokens=1000
            )

        result_text = response.choices[0].message.content
        os.remove(filepath)  # Bersihkan file setelah selesai
        return jsonify({"result": result_text})
    except Exception as e:
        return jsonify({"error": f"Gagal menganalisis gambar: {str(e)}"})
    


@app.route("/voicevox_tts", methods=["POST"])
def voicevox_tts():
    data = request.get_json()
    text = data.get("text", "")
    speaker = int(data.get("speaker", 1))

    # Buat query audio
    query_res = requests.post(
        f"http://localhost:50021/audio_query?text={text}&speaker={speaker}"
    )
    query_res.raise_for_status()
    audio_query = query_res.json()

    # Synthesis audio
    synth_res = requests.post(
        f"http://localhost:50021/synthesis?speaker={speaker}",
        headers={"Content-Type": "application/json"},
        data=json.dumps(audio_query)
    )
    synth_res.raise_for_status()

    # Return audio as WAV
    return send_file(
        io.BytesIO(synth_res.content),
        mimetype="audio/wav",
        as_attachment=False
    )
    
@app.route('/api/text-to-speech', methods=['POST'])
def text_to_speech():
    data = request.get_json()
    text = data.get("text", "")
    if not text:
        return jsonify({"error": "No text provided"}), 400
    
    # Logic TTS-mu di sini (bisa pakai pyttsx3 atau gTTS atau lainnya)
    # Simulasi respon dummy:
    return jsonify({"audio_url": "/static/tts_output.mp3"})


@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        text = request.form["text"]
        speaker = int(request.form.get("speaker", 1))

        # 1. Buat audio query
        query_url = f"{VOICEVOX_ENGINE_URL}/audio_query?text={text}&speaker={speaker}"
        query_response = requests.post(query_url)
        audio_query = query_response.json()

        # 2. Synthesize audio
        synth_url = f"{VOICEVOX_ENGINE_URL}/synthesis?speaker={speaker}"
        synth_response = requests.post(
            synth_url,
            headers={"Content-Type": "application/json"},
            data=json.dumps(audio_query)
        )

        # 3. Kirim sebagai file .wav
        return send_file(
            io.BytesIO(synth_response.content),
            mimetype="audio/wav",
            as_attachment=True,
            download_name="output.wav"
        )
        

    return render_template("chatbot.html")
# Home route


# Main entry point
if __name__ == '__main__':
    app.run(debug=True)
